"""
PPT Translation API Router
Handles PowerPoint upload, translation, and download with real-time progress
"""

from fastapi import APIRouter, UploadFile, File, Form, Depends, HTTPException
from fastapi.responses import StreamingResponse, FileResponse
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy.orm.attributes import flag_modified
from sqlalchemy import select
from typing import List, Optional
import json
import copy as copy_module
from datetime import datetime
import asyncio

from models.database import Job, get_db, AsyncSessionLocal
from services.ppt_service import create_translated_ppt, has_image_on_left, detect_language_from_presentation
from utils.helpers import save_job_file, get_api_key
from config.settings import settings
from pptx import Presentation
import io

router = APIRouter(prefix="/api/ppt", tags=["ppt"])

# Store background tasks to prevent garbage collection
background_tasks = set()

@router.post("/preview")
async def preview_slides(
    file: UploadFile = File(...),
):
    """Preview slides and detect which ones have images on the left"""
    try:
        # Read file content
        file_content = await file.read()
        
        # Verify it's a valid PPTX
        if not file.filename or not file.filename.endswith('.pptx'):
            raise HTTPException(status_code=400, detail="File must be a .pptx file")
        
        # Load presentation
        try:
            prs = Presentation(io.BytesIO(file_content))
        except Exception as pptx_error:
            raise HTTPException(status_code=400, detail=f"Invalid PowerPoint file: {str(pptx_error)}")
        
        slides_info = []
        for idx, slide in enumerate(prs.slides):
            try:
                has_img, image = has_image_on_left(slide)
                
                slide_info = {
                    "index": idx,
                    "has_image": has_img,
                    "selected": has_img,  # Default: select slides with images
                    "layout": slide.slide_layout.name if hasattr(slide, 'slide_layout') else "Unknown"
                }
                slides_info.append(slide_info)
            except Exception as slide_error:
                # If individual slide fails, still include it
                slide_info = {
                    "index": idx,
                    "has_image": False,
                    "selected": False,
                    "layout": "Error",
                    "error": str(slide_error)
                }
                slides_info.append(slide_info)
        
        # Detect language from selected slides (those with images)
        selected_indices = [s["index"] for s in slides_info if s["selected"]]
        detected_lang = None
        if selected_indices:
            detected_lang = detect_language_from_presentation(prs, selected_indices)
            print(f"🌍 Detected language: {detected_lang}")
        
        return {
            "total_slides": len(prs.slides),
            "slides": slides_info,
            "detected_lang": detected_lang
        }
    
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Error previewing slides: {str(e)}")


@router.post("/translate")
async def translate_ppt(
    file: UploadFile = File(...),
    provider: str = Form(...),
    source_lang: str = Form("es"),
    target_lang: str = Form("en"),
    selected_slides: str = Form("[]"),  # JSON string of slide indices
    base_font_size: int = Form(11),
    title_size_adjustment: int = Form(5),
    subject_size_adjustment: int = Form(1),
    preserve_colors: bool = Form(True),
    db: AsyncSession = Depends(get_db)
):
    """Start PPT translation job"""
    try:
        # Parse selected slides
        selected_slide_indices = json.loads(selected_slides)
        
        # Read file
        file_content = await file.read()
        
        # Get API key
        api_key = await get_api_key(provider, db)
        claude_api_key = await get_api_key("claude", db)
        
        if not api_key and provider in ["claude", "openrouter"]:
            raise HTTPException(
                status_code=400,
                detail=f"API key for {provider} not configured. Please set it in admin panel."
            )
        
        # Get model
        if provider == "claude":
            model = settings.default_claude_model
        elif provider in ["openrouter", "ocr_free"]:
            model = settings.default_openrouter_model
        else:
            model = "offline"
        
        # Create job record
        job = Job(
            job_type="ppt_translation",
            status="processing",
            input_filename=file.filename,
            provider=provider,
            source_lang=source_lang,
            target_lang=target_lang,
            settings_used={
                "base_font_size": base_font_size,
                "title_size_adjustment": title_size_adjustment,
                "subject_size_adjustment": subject_size_adjustment,
                "preserve_colors": preserve_colors,
                "selected_slides": selected_slide_indices,
                "selected_slides_count": len(selected_slide_indices)  # For progress calculation
            }
        )
        
        db.add(job)
        await db.commit()
        await db.refresh(job)
        
        # Create and store task reference
        task = asyncio.create_task(
            process_translation_background(
                job.id, file_content, selected_slide_indices, 
                source_lang, target_lang, provider, api_key, model,
                base_font_size, title_size_adjustment, subject_size_adjustment,
                claude_api_key
            )
        )

        # Store task reference to prevent garbage collection
        background_tasks.add(task)
        task.add_done_callback(background_tasks.discard)

        print(f"✅ Background task started for job {job.id}")
        
        # Return job ID for progress tracking
        return {
            "job_id": job.id,
            "status": "processing",
            "message": "Translation started"
        }
    
    except Exception as e:
        import traceback
        print("❌ Error starting translation:")
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


async def process_translation_background(
    job_id: int,
    file_content: bytes,
    selected_slides: list,
    source_lang: str,
    target_lang: str,
    provider: str,
    api_key: str,
    model: str,
    base_font_size: int,
    title_adjustment: int,
    subject_adjustment: int,
    claude_api_key: Optional[str]
):
    """Background task to process translation"""
    print(f"🎯 INSIDE background task for job {job_id}")
    print(f"Processing {len(selected_slides)} slides with {provider}")
    
    from sqlalchemy import select
    
    output_path = None
    async with AsyncSessionLocal() as db:
        try:
            # Get job
            print(f"📊 Getting job {job_id} from database...")
            result = await db.execute(select(Job).where(Job.id == job_id))
            job = result.scalar_one_or_none()
            
            if not job:
                print(f"❌ Job {job_id} not found!")
                return
            print(f"✅ Job found: {job.input_filename}")
            
            # Load presentation
            print(f"📄 Loading presentation...")
            prs = Presentation(io.BytesIO(file_content))
            print(f"✅ Presentation loaded: {len(prs.slides)} total slides")

            # Detect slides with images for accurate progress
            slides_with_images = []
            for idx, slide in enumerate(prs.slides):
                if idx not in selected_slides:
                    continue
                try:
                    has_img, _ = has_image_on_left(slide)
                    if has_img:
                        slides_with_images.append(idx)
                except Exception:
                    continue

            steps_per_slide = 3
            total_steps = len(slides_with_images) * steps_per_slide

            job.status = "detecting"
            job.slides_processed = 0
            job.total_slides = total_steps
            if isinstance(job.settings_used, dict):
                job.settings_used["slides_detected"] = len(slides_with_images)
                job.settings_used["progress_steps_per_slide"] = steps_per_slide
            await db.commit()
            
            # Process translation
            settings_dict = {
                "base_font_size": base_font_size,
                "title_size_adjustment": title_adjustment,
                "subject_size_adjustment": subject_adjustment,
                "progress_steps_per_slide": steps_per_slide,
                "enable_offline_fallback": settings.enable_offline_fallback
            }
            if claude_api_key:
                settings_dict["claude_api_key"] = claude_api_key
            
            # Progress callback
            async def update_progress(current, total, status, data):
                slide_index = data.get("slide_index") if data else None
                print(f"📈 Progress: {current}/{total} - {status} - Slide #{slide_index}" if slide_index else f"📈 Progress: {current}/{total} - {status}")
                result = await db.execute(select(Job).where(Job.id == job_id))
                job = result.scalar_one_or_none()
                if job:
                    # Refresh to get latest data from DB
                    await db.refresh(job)
                    
                    # Check if job was cancelled
                    cancelled_flag = job.settings_used.get("cancelled") if isinstance(job.settings_used, dict) else False
                    print(f"🔍 Checking cancellation: cancelled={cancelled_flag}, settings_used type={type(job.settings_used)}")
                    
                    if cancelled_flag:
                        print(f"🛑🛑🛑 CANCELLATION DETECTED! Job {job_id} was cancelled, raising exception NOW!")
                        raise Exception("Translation cancelled by user")
                    
                    job.slides_processed = current
                    
                    # Store current slide index for display
                    if slide_index and isinstance(job.settings_used, dict):
                        job.settings_used["current_slide_index"] = slide_index
                        flag_modified(job, "settings_used")
                    
                    # Don't overwrite job.total_slides - it's the total PPT slides count
                    # The 'total' parameter here is actually selected_slides_count
                    if status in ["processing", "completed", "failed"]:
                        job.status = "processing" if status != "failed" else "failed"
                    await db.commit()
            
            # Do translation
            print(f"🚀 Calling create_translated_ppt...")
            output_bytes, stats = await create_translated_ppt(
                file_content,
                selected_slides,
                source_lang,
                target_lang,
                provider,
                api_key,
                model,
                settings_dict,
                update_progress
            )
            
            print(f"✅ Translation completed! Stats: {stats}")
            print(f"⏱️  Elapsed seconds from stats: {stats.get('elapsed_seconds', 'NOT FOUND')}")
            
            # Save output file
            print(f"💾 Saving output file...")
            ts = datetime.utcnow().strftime("%y%m%d_%H%M")
            stem = job.input_filename.rsplit(".", 1)[0] if "." in job.input_filename else job.input_filename
            output_filename = f"translated_{stem}_{ts}.pptx"
            output_path = save_job_file(output_bytes, output_filename, "outputs")
            print(f"✅ File saved to: {output_path}")

            # Update job
            result = await db.execute(select(Job).where(Job.id == job_id))
            job = result.scalar_one_or_none()
            if job:
                job.status = "completed"
                job.output_path = str(output_path)
                job.output_filename = output_filename
                job.slides_processed = stats["processed_slides"]
                job.total_slides = stats["total_slides"]
                job.completed_at = datetime.utcnow()
                
                # Save all stats in settings_used for frontend access
                # Create a new dict to ensure SQLAlchemy detects the change
                updated_settings = dict(job.settings_used) if isinstance(job.settings_used, dict) else {}
                updated_settings.update({
                    "elapsed_seconds": stats.get("elapsed_seconds", 0),
                    "total_cost": stats.get("total_cost", 0),
                    "total_input_tokens": stats.get("total_input_tokens", 0),
                    "total_output_tokens": stats.get("total_output_tokens", 0),
                    "slide_methods": stats.get("slide_methods", []),
                    "failed_slides": stats.get("failed_slides", 0)
                })
                job.settings_used = updated_settings  # Assign new object
                print(f"💾 Saved to settings_used: elapsed_seconds={job.settings_used.get('elapsed_seconds')}, total_cost={job.settings_used.get('total_cost')}")
                
                # Estimate cost (Claude)
                if provider == "claude":
                    job.estimated_cost = stats["processed_slides"] * 0.003
                
                await db.commit()
                print(f"🎉 Job {job_id} marked as completed!")
        
        except Exception as e:
            error_msg = str(e)
            is_cancellation = "cancelled by user" in error_msg.lower()
            
            if is_cancellation:
                print(f"🛑 Translation cancelled by user for job {job_id}")
            else:
                print(f"❌❌❌ ERROR in background task: {e}")
                import traceback
                traceback.print_exc()
            
            # Update job with error
            result = await db.execute(select(Job).where(Job.id == job_id))
            job = result.scalar_one_or_none()
            if job:
                job.status = "failed"
                job.error_message = "Cancelled by user" if is_cancellation else error_msg
                job.completed_at = datetime.utcnow()
                await db.commit()
            
            if not is_cancellation:
                print(f"Translation error: {e}")



@router.get("/progress/{job_id}")
async def get_progress(
    job_id: int,
    db: AsyncSession = Depends(get_db)
):
    """Get job progress via Server-Sent Events"""
    
    async def event_generator():
        """Generate SSE events for job progress"""
        max_iterations = 300  # 5 minutes max
        iteration = 0
        
        while iteration < max_iterations:
            # Fetch job from database
            from sqlalchemy import select
            result = await db.execute(select(Job).where(Job.id == job_id))
            job = result.scalar_one_or_none()
            
            if not job:
                yield f"data: {json.dumps({'error': 'Job not found'})}\n\n"
                break
            
            # Send progress update
            progress_data = {
                "status": job.status,
                "slides_processed": job.slides_processed,
                "total_slides": job.total_slides,
                "error_message": job.error_message
            }

            # Add friendly message and detection info when available
            steps_per_slide = 1
            slides_detected = None
            if isinstance(job.settings_used, dict):
                steps_per_slide = job.settings_used.get("progress_steps_per_slide", 1)
                slides_detected = job.settings_used.get("slides_detected")

            if job.status == "detecting":
                if slides_detected is not None:
                    progress_data["message"] = f"Slides detected: {slides_detected}"
                else:
                    progress_data["message"] = "Detecting slides with images..."
            elif job.status in ["processing", "completed"]:
                if slides_detected is not None and steps_per_slide > 0:
                    slide_num = min(slides_detected, (job.slides_processed // steps_per_slide) + 1) if slides_detected > 0 else 0
                    progress_data["message"] = f"Translating slide {slide_num}/{slides_detected}"
                    progress_data["slides_detected"] = slides_detected
                else:
                    progress_data["message"] = "Translating slides..."
            
            yield f"data: {json.dumps(progress_data)}\n\n"
            
            # If job is completed or failed, stop
            if job.status in ["completed", "failed"]:
                break
            
            await asyncio.sleep(1)
            iteration += 1
        
        # Close connection
        yield f"data: {json.dumps({'status': 'done'})}\n\n"
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no"
        }
    )


@router.get("/download/{job_id}")
async def download_translation(
    job_id: int,
    db: AsyncSession = Depends(get_db)
):
    """Download translated PPT file"""
    try:
        print(f"🔍 PPT Download request for job_id: {job_id}")
        
        from sqlalchemy import select
        result = await db.execute(select(Job).where(Job.id == job_id))
        job = result.scalar_one_or_none()
        
        if not job:
            print(f"❌ Job {job_id} not found in database")
            raise HTTPException(status_code=404, detail="Job not found")
        
        print(f"✅ Job found: status={job.status}")
        print(f"📄 output_path from DB: {job.output_path}")
        print(f"📄 output_filename from DB: {job.output_filename}")
        
        if job.status != "completed":
            print(f"❌ Job not completed, status: {job.status}")
            raise HTTPException(status_code=400, detail="Job not completed yet")
        
        from pathlib import Path
        import os
        
        file_path = Path(job.output_path)
        print(f"🔍 Looking for file at: {file_path}")
        print(f"📁 File exists? {file_path.exists()}")
        print(f"📂 Settings output_dir: {settings.output_dir}")
        
        # List files in output directory
        try:
            if settings.output_dir.exists():
                files_in_dir = list(settings.output_dir.iterdir())
                print(f"📂 Files in {settings.output_dir}:")
                for f in files_in_dir:
                    print(f"   - {f.name}")
            else:
                print(f"❌ Output directory doesn't exist: {settings.output_dir}")
        except Exception as e:
            print(f"⚠️ Error listing directory: {e}")
        
        if not file_path.exists():
            print(f"❌ File not found at path: {file_path}")
            
            # Try to find it by filename only
            potential_path = settings.output_dir / job.output_filename
            print(f"🔍 Trying alternate path: {potential_path}")
            print(f"📁 Alternate exists? {potential_path.exists()}")
            
            if potential_path.exists():
                print(f"✅ Found at alternate location! Using that.")
                file_path = potential_path
            else:
                raise HTTPException(status_code=404, detail=f"File not found: {file_path}")
        
        print(f"✅ File found! Size: {file_path.stat().st_size} bytes")
        print(f"📤 Returning file: {job.output_filename}")
        
        return FileResponse(
            file_path,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=job.output_filename
        )
    
    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Error in PPT download: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/history")
async def get_history(
    limit: int = 20,
    db: AsyncSession = Depends(get_db)
):
    """Get translation history"""
    from sqlalchemy import select
    result = await db.execute(
        select(Job)
        .where(Job.job_type == "ppt_translation")
        .order_by(Job.created_at.desc())
        .limit(limit)
    )
    jobs = result.scalars().all()
    
    return {
        "jobs": [job.to_dict() for job in jobs]
    }


@router.post("/merge")
async def merge_translations(
    job_id_1: int = Form(...),
    job_id_2: int = Form(...),
    db: AsyncSession = Depends(get_db)
):
    """Merge two translation job outputs into a single ordered PPTX (original slide order)."""
    try:
        # Load both jobs
        result1 = await db.execute(select(Job).where(Job.id == job_id_1))
        job1 = result1.scalar_one_or_none()
        result2 = await db.execute(select(Job).where(Job.id == job_id_2))
        job2 = result2.scalar_one_or_none()

        if not job1 or not job2:
            raise HTTPException(status_code=404, detail="One or both jobs not found")
        if job1.status != "completed" or job2.status != "completed":
            raise HTTPException(status_code=400, detail="Both jobs must be completed before merging")

        # Resolve output file paths
        from pathlib import Path
        def resolve_path(job):
            p = Path(job.output_path)
            if not p.exists():
                p = settings.output_dir / job.output_filename
            if not p.exists():
                raise HTTPException(status_code=404, detail=f"Output file not found for job {job.id}")
            return p

        path1 = resolve_path(job1)
        path2 = resolve_path(job2)

        # Build (original_1based_idx → output_0based_position) maps from slide_methods
        def successful_slide_map(slide_methods):
            """Returns [(original_1based_idx, output_0based_pos), ...] sorted by original index."""
            ok = [m for m in slide_methods if m.get("method") and m.get("method") != "unknown"]
            ok.sort(key=lambda m: m["slide"])
            return [(m["slide"], i) for i, m in enumerate(ok)]

        methods1 = (job1.settings_used or {}).get("slide_methods", [])
        methods2 = (job2.settings_used or {}).get("slide_methods", [])

        slides1 = successful_slide_map(methods1)  # from job1 output
        slides2 = successful_slide_map(methods2)  # from job2 output (retry)

        # Combine: job2 wins on duplicate original indices (retry takes priority)
        seen = {}
        for orig_idx, out_pos in slides1:
            seen[orig_idx] = (1, out_pos)
        for orig_idx, out_pos in slides2:
            seen[orig_idx] = (2, out_pos)  # retry overrides

        combined = sorted(seen.items())  # [(orig_idx, (job_num, out_pos)), ...]
        print(f"🔀 Merge: job1={len(slides1)} slides, job2={len(slides2)} slides → merged={len(combined)}")

        prs1 = Presentation(str(path1))
        prs2 = Presentation(str(path2))

        merged = Presentation()
        merged.slide_width = prs1.slide_width
        merged.slide_height = prs1.slide_height

        for orig_idx, (job_num, out_pos) in combined:
            src_prs = prs1 if job_num == 1 else prs2
            src_slide = src_prs.slides[out_pos]
            dest_slide = merged.slides.add_slide(merged.slide_layouts[6])  # blank

            for shape in src_slide.shapes:
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    dest_slide.shapes.add_picture(
                        io.BytesIO(shape.image.blob),
                        shape.left, shape.top, shape.width, shape.height
                    )
                else:
                    dest_slide.shapes._spTree.append(copy_module.deepcopy(shape.element))

            print(f"  ✅ Slide orig#{orig_idx} from job{job_num}[{out_pos}]")

        # Save merged output
        timestamp = datetime.now().strftime('%Y-%m-%d_%Hh%M')
        filename = f"Merged_Translation_{timestamp}.pptx"
        buf = io.BytesIO()
        merged.save(buf)
        buf.seek(0)

        file_path = save_job_file(buf.getvalue(), filename, "outputs")

        merge_job = Job(
            job_type="ppt_merge",
            status="completed",
            input_filename=f"merge_of_job{job_id_1}_and_job{job_id_2}",
            output_filename=filename,
            output_path=str(file_path),
            completed_at=datetime.now(),
            settings_used={"merged_from": [job_id_1, job_id_2], "total_slides": len(combined)}
        )
        db.add(merge_job)
        await db.commit()
        await db.refresh(merge_job)

        print(f"✅ Merge complete: {filename} ({len(combined)} slides), job_id={merge_job.id}")
        return {"job_id": merge_job.id, "total_slides": len(combined), "filename": filename}

    except HTTPException:
        raise
    except Exception as e:
        print(f"❌ Merge error: {e}")
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


def process_translation_background_sync(
    job_id: int,
    file_content: bytes,
    selected_slides: list,
    source_lang: str,
    target_lang: str,
    provider: str,
    api_key: str,
    model: str,
    base_font_size: int,
    title_adjustment: int,
    subject_adjustment: int,
    claude_api_key: Optional[str] = None
):
    """Synchronous wrapper for background processing"""
    print(f"🔥 Background task ACTUALLY started for job {job_id}")
    
    # Run the async function in a new event loop
    import asyncio
    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)
    
    try:
        loop.run_until_complete(
            process_translation_background(
                job_id, file_content, selected_slides,
                source_lang, target_lang, provider, api_key, model,
                base_font_size, title_adjustment, subject_adjustment,
                claude_api_key
            )
        )
    finally:
        loop.close()
    
    print(f"✅ Background task completed for job {job_id}")


@router.post("/cancel/{job_id}")
async def cancel_translation(
    job_id: int,
    db: AsyncSession = Depends(get_db)
):
    """Cancel an ongoing translation job"""
    result = await db.execute(select(Job).where(Job.id == job_id))
    job = result.scalar_one_or_none()
    
    if not job:
        raise HTTPException(status_code=404, detail="Job not found")
    
    if job.status not in ["pending", "processing"]:
        raise HTTPException(status_code=400, detail="Job cannot be cancelled (already completed or failed)")
    
    # Mark job as cancelled with a flag the background process can check
    job.status = "failed"
    job.error_message = "Cancelled by user"
    job.completed_at = datetime.utcnow()
    
    # Set cancellation flag in settings_used so background task can detect it
    if isinstance(job.settings_used, dict):
        job.settings_used["cancelled"] = True
        flag_modified(job, "settings_used")
        print(f"🚫 Setting cancelled=True for job {job_id}, settings_used={job.settings_used}")
    else:
        print(f"⚠️  Cannot set cancelled flag - settings_used is not a dict: {type(job.settings_used)}")
    
    await db.commit()
    print(f"✅ Cancel committed to DB for job {job_id}")
    
    return {"message": "Translation cancelled successfully"}
