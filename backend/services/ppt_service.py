"""
PowerPoint Translation Service
Handles OCR, LLM translation, and PPT generation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from PIL import Image, ImageEnhance, ImageFilter
from langdetect import detect, LangDetectException
import io
import json
import base64
import time
import anthropic
import requests
import pytesseract
from typing import Dict, List, Optional, Tuple
import asyncio
from concurrent.futures import ThreadPoolExecutor

# Thread pool for blocking operations
executor = ThreadPoolExecutor(max_workers=4)


def encode_image_to_base64(image: Image.Image) -> str:
    """Encode PIL image to base64"""
    try:
        width, height = image.size
        image = image.resize((int(width * 2), int(height * 2)), Image.LANCZOS)
        # Improve OCR readability (contrast + sharpness)
        image = ImageEnhance.Contrast(image).enhance(1.2)
        image = ImageEnhance.Sharpness(image).enhance(1.3)
    except Exception:
        pass
    buffered = io.BytesIO()
    image.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode()


def _create_compact_vision_prompt(source_lang: str, target_lang: str) -> str:
    """Compact vision prompt for retry when full prompt produces truncated JSON."""
    lang_names = {"fr": "French", "es": "Spanish", "it": "Italian", "en": "English"}
    return f"""Extract all text from this image and translate from {lang_names[source_lang]} to {lang_names[target_lang]}.

Return ONLY a JSON object. Be very concise — omit whitespace in JSON.
For each line of text, create one entry. Use empty blocks array for blank lines.
Text color: look only at letter color (ignore underlines). Bold: true if thick strokes.

Format:
{{"lines":[{{"alignment":"left","blocks":[{{"text":"translated text","bold":false,"color":"black","size":"normal","alignment":"left"}}]}},{{"alignment":"left","blocks":[]}}]}}

Rules:
- End each text block with \n only if there is a real line break inside a group
- Multiple bullets/items in the same visual group go as separate blocks in the same line
- Translate faithfully, keep product names unchanged
- Output ONLY valid JSON, nothing else"""


def create_vision_prompt(source_lang: str, target_lang: str) -> str:
    """Create prompt for Claude to extract and translate text"""
    lang_names = {"fr": "French", "es": "Spanish", "it": "Italian", "en": "English"}
    
    prompt = f"""You are a precise OCR and translation system.

**YOUR TASK:**
1. Extract ALL text from the image
2. Translate from {lang_names[source_lang]} to {lang_names[target_lang]}
3. Copy EXACTLY the same formatting for each text element

**FORMATTING TO COPY EXACTLY:**

For EACH piece of text, observe and copy:

1. **TEXT COLOR** (CRITICAL - Read this carefully!):
   
   🚨 **CRITICAL WARNING - COLORED UNDERLINES:**
   Many words in this document have COLORED underlines but different text colors!
   - "ballonnements" = BLACK text + RED underline → color: "black"
   - "gonfiement" = BLUE text + RED underline → color: "blue"  
   - "digestion" = BLACK text + RED underline → color: "black"
   - "CONSEILLÉ" = GREEN text + BLUE underline → color: "green"
   - Some word = BLACK text + BLUE underline → color: "black"
   DO NOT confuse the underline color with the text color!
   
   ⚠️ **HOW TO IDENTIFY TEXT COLOR:**
   - Look ONLY at the color of the letters/characters themselves
   - Completely IGNORE any underline, highlight, or background color
   
   ⚠️ **COMMON MISTAKE TO AVOID:**
   If you see BLUE text with a RED underline:
   ❌ WRONG: color: "red" (you looked at underline)
   ✅ CORRECT: color: "blue" (you looked at the text itself)
   
   **Visual Test:**
   Ask yourself: "If I remove all underlines and backgrounds, what color are the letters?"
   That's the color to report.
      
2. **BOLD** (is the text thick/heavy or normal weight?):
   - true = letters are thick/bold/heavy
   - false = letters are normal weight
   
   ⚠️ **HOW TO DETECT BOLD:**
   - Compare the stroke/line thickness of the letters
   - Bold letters have thicker strokes than regular letters
   - Look at the letter edges - are they thicker/heavier?
   
   ⚠️ **COMMON MISTAKE:**
   Titles are OFTEN bold even if they're colored!
   - Blue title in bold → bold: true, color: "blue"
   - Green title in bold → bold: true, color: "green"
   
   **Bold Test:**
   Compare this text to regular body text:
   - Are the letter strokes thicker? → bold: true
   - Same thickness as body text? → bold: false
   
   **Important:** Color does NOT affect bold detection
   - Blue text can be bold OR not bold
   - Green text can be bold OR not bold
   - Check the stroke thickness independently of color
   
3. **SIZE** (compare to normal body text):
   - "large" = bigger than body text (titles, headers)
   - "normal" = regular body text size
   - "small" = smaller than body text (footnotes)
   
4. **LINE BREAKS** (IMPORTANT):
   - Add "\\n" ONLY when there is an ACTUAL line break in the original
   - Keep text together on the same line if it appears together
   - DO NOT split titles or sentences unnecessarily
   
   ⚠️ **COMMON MISTAKE:**
   If you see "14- Supplement Digestion (Single Scale)" on ONE line:
   ❌ WRONG: Split into multiple blocks with \\n between words
   ✅ CORRECT: Keep it as one block without \\n inside
   
   **Line Break Test:**
   Ask: "Does the next word start on a NEW line below, or continue on the SAME line?"
   - Same line → NO \\n
   - New line below → ADD \\n at the end
   
5. **ALIGNMENT** (IMPORTANT):
   - "left" = text starts at the left margin (most common)
   - "center" = text is centered horizontally
   - "right" = text starts at the right margin
   
   ⚠️ **COMMON MISTAKE:**
   Body text paragraphs are almost ALWAYS "left" aligned
   Only titles/headers are sometimes centered
   
   **Alignment Test:**
   - Look at where the text starts horizontally
   - If it starts at the left edge → "left"
   - If it's in the middle of the page → "center"
   - Most body text = "left"

**CRITICAL RULES:**

❌ DO NOT add formatting that isn't there
❌ DO NOT make assumptions about product names or keywords
❌ DO NOT let underline colors influence text color
❌ DO NOT assume certain words should be bold or colored
❌ DO NOT split text onto multiple lines if it appears on ONE line in the image
❌ DO NOT center text that is left-aligned in the image
✅ ONLY copy what you actually SEE in the letters themselves
✅ Ignore all underlines, highlights, and backgrounds when determining text color
✅ If text looks black to you (ignoring underline), write "black"
✅ If text is NOT bold (regardless of what word it is), write "bold": false
✅ Keep titles together on one line if they appear on one line
✅ Most body text is "left" aligned - only special titles are centered
✅ Split into separate blocks whenever formatting changes

**OUTPUT FORMAT (PREFERRED):**

Return JSON with explicit lines to preserve line breaks exactly.
Each line contains blocks. Empty lines use an empty blocks array.

{{{{
  "lines": [
    {{{{
      "alignment": "left",
      "blocks": [
        {{{{"text": "Some text ", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}},
        {{{{"text": "bold word", "bold": true, "color": "black", "size": "normal", "alignment": "left"}}}}
      ]
    }}}},
    {{{{
      "alignment": "left",
      "blocks": []
    }}}},
    {{{{
      "alignment": "center",
      "blocks": [
        {{{{"text": "Newsletter", "bold": true, "color": "black", "size": "large", "alignment": "center"}}}}
      ]
    }}}}
  ]
}}}}

If you cannot produce "lines", return the older "blocks" format.

**EXAMPLES WITH UNDERLINES:**

Example 1: Blue text "Digestion" with red underline underneath
Visual: The word "Digestion" appears in blue letters, with a red line under it
→ {{{{"text": "Digestion", "bold": false, "color": "blue", "size": "normal", "alignment": "left"}}}}
(Color is BLUE because the LETTERS are blue. The red underline is decoration, not text color)

Example 1b: Blue title "14- Health Description (Option)" in BOLD
Visual: Blue letters that are thick/heavy (bold weight)
→ {{{{"text": "14- Health Description (Option)\\n", "bold": true, "color": "blue", "size": "large", "alignment": "left"}}}}
(Color is BLUE and bold is TRUE because letters are both blue AND thick)

Example 1c: Green title "ColigasFast: your natural ally" in BOLD
Visual: Green letters that are thick/heavy (bold weight)
→ {{{{"text": "ColigasFast: your natural ally\\n", "bold": true, "color": "green", "size": "large", "alignment": "left"}}}}
(Color is GREEN and bold is TRUE because letters are both green AND thick)

Example 2: BLACK text "ballonnements" with RED underline (CRITICAL!)
Visual: The word "ballonnements" in BLACK letters, RED straight line underneath
→ {{{{"text": "ballonnements", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}}
(Color is BLACK because the LETTERS are black. IGNORE the red underline completely!)

Example 2b: BLACK text "gonfiement" with RED squiggly underline
Visual: The word "gonfiement" in BLACK letters, RED wavy line underneath  
→ {{{{"text": "gonfiement", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}}
(Color is BLACK because the LETTERS are black. The red underline is NOT the text color!)

Example 2c: BLACK text "CONSEILLÉ" with BLUE underline (CRITICAL!)
Visual: The word "CONSEILLÉ" in BLACK letters, BLUE straight line underneath
→ {{{{"text": "CONSEILLÉ", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}}
(Color is BLACK because the LETTERS are black. IGNORE the blue underline completely!)

Example 3: Black text "bloating" with red squiggly underline (spell-check style)
Visual: The word "bloating" in black letters, red squiggly line underneath
→ {{{{"text": "bloating", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}}
(Color is BLACK because the LETTERS are black, ignore the red underline)

Example 4: Title "14- Supplement Digestion (Single Scale)" all in blue
Visual: All text in this title is blue letters ON ONE LINE
→ {{{{"text": "14- Supplement Digestion (Single Scale)\\n", "bold": true, "color": "blue", "size": "large", "alignment": "left"}}}}
(Everything together because it's on ONE line. \\n only at the END for line break after title)

Example 5: Body paragraph (multiple lines of left-aligned text)
Visual: Text starts at left margin, flows naturally to next line
→ [
  {{{{"text": "This is a sentence that continues ", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}},
  {{{{"text": "on the same line.\\n", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}},
  {{{{"text": "This is a new line below.", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}}
]
(Each line break in original → \\n at end. All "alignment": "left" for body text)

Example 6: Centered title
Visual: "Newsletter" appears centered horizontally on the page
→ {{{{"text": "Newsletter\\n", "bold": true, "color": "black", "size": "large", "alignment": "center"}}}}
(alignment": "center" because it's visually centered)

Example 7: Bold question in BLACK
Visual: "Avec qui vivez-vous ?" in BLACK letters that are thick/heavy (bold weight)
→ {{{{"text": "Avec qui vivez-vous ?\\n", "bold": true, "color": "black", "size": "normal", "alignment": "left"}}}}
(Bold is TRUE because the letters are thick, even though it's a regular question, not a title)

Example 8: "Aboca" in regular weight black text
→ {{{{"text": "Aboca", "bold": false, "color": "black", "size": "normal", "alignment": "left"}}}}
(NOT bold because it doesn't LOOK bold)

Example 9: "ColiGas Fast" in bold black text
→ {{{{"text": "ColiGas Fast", "bold": true, "color": "black", "size": "normal", "alignment": "left"}}}}
(Bold because the letters are visually heavier/thicker)

**FINAL REMINDERS:**
- Text color = letter color ONLY (ignore underlines completely)
- Bold = visual stroke thickness ONLY (ignore content/meaning)
- **IMPORTANT**: Colored text can ALSO be bold! Check both independently:
  - Blue text with thick strokes → bold: true, color: "blue"
  - Green text with thick strokes → bold: true, color: "green"
  - Red text with normal strokes → bold: false, color: "red"
- Translate text to {lang_names[target_lang]}
- Copy formatting EXACTLY as you see it in the letters
- Don't assume or add formatting
- Return ONLY valid JSON, no markdown blocks"""
    
    return prompt



SIZE_MAP = {
    "huge": 7,
    "very-large": 5,
    "large": 3,
    "normal": 0,
    "small": -2,
    "tiny": -4
}


def _parse_vision_response(result: dict) -> Optional[Dict]:
    """Parse JSON response from Claude or OpenRouter vision into PPT structure.
    Handles both 'lines' format (preferred) and legacy 'blocks' format.
    """
    if "lines" in result:
        structure = []
        for line in result.get("lines", []):
            line_blocks = []
            for block in line.get("blocks", []):
                size_relative = SIZE_MAP.get(block.get("size", "normal"), 0)
                line_blocks.append({
                    "type": "paragraph",
                    "original": block.get("text", ""),
                    "translated": block.get("text", ""),
                    "formatting": {
                        "bold": block.get("bold", False),
                        "size_relative": size_relative,
                        "color": block.get("color", "black"),
                        "alignment": block.get("alignment", line.get("alignment", "left"))
                    },
                    "has_emoji": True
                })
            structure.append({
                "line_blocks": line_blocks,
                "alignment": line.get("alignment", "left")
            })
        return {"original_text": "", "translated_text": "", "structure": structure}

    if "blocks" in result:
        structure = []
        for block in result["blocks"]:
            size_relative = SIZE_MAP.get(block.get("size", "normal"), 0)
            structure.append({
                "type": "paragraph",
                "original": block.get("text", ""),
                "translated": block.get("text", ""),
                "formatting": {
                    "bold": block.get("bold", False),
                    "size_relative": size_relative,
                    "color": block.get("color", "black"),
                    "alignment": block.get("alignment", "left")
                },
                "has_emoji": True
            })
        return {"original_text": "", "translated_text": "", "structure": structure}

    return None

def _repair_truncated_json(text: str) -> Optional[str]:
    """Attempt to repair a truncated JSON string by closing open brackets/braces.
    Returns repaired JSON string, or None if unrepairable."""
    text = text.strip()
    if not text:
        return None

    # Remove trailing partial token (incomplete string, key, or value)
    # Find last complete value boundary
    # Walk backwards to find last clean delimiter
    for i in range(len(text) - 1, -1, -1):
        c = text[i]
        if c in ('}', ']', '"', '0123456789', 'e', 'l', 's', 'n', 'r', 'u'):
            # Try to close from here
            candidate = text[:i + 1]
            break
    else:
        return None

    # Count open brackets/braces to determine what needs closing
    stack = []
    in_string = False
    escape_next = False
    for ch in candidate:
        if escape_next:
            escape_next = False
            continue
        if ch == '\\':
            escape_next = True
            continue
        if ch == '"' and not in_string:
            in_string = True
        elif ch == '"' and in_string:
            in_string = False
        elif not in_string:
            if ch in ('{', '['):
                stack.append(ch)
            elif ch == '}':
                if stack and stack[-1] == '{':
                    stack.pop()
            elif ch == ']':
                if stack and stack[-1] == '[':
                    stack.pop()

    if not stack:
        return candidate  # Already balanced

    # Close open brackets in reverse order
    closing = {'[': ']', '{': '}'}
    closers = ''.join(closing[c] for c in reversed(stack))

    # Make sure we end on a clean boundary before closing
    # Strip trailing comma or incomplete content
    candidate = candidate.rstrip().rstrip(',').rstrip()
    repaired = candidate + closers
    return repaired


async def call_claude_vision(
    image: Image.Image,
    prompt: str,
    api_key: str,
    model: str = "claude-sonnet-4-20250514"
) -> Optional[Dict]:
    """Call Claude API for vision analysis"""
    try:
        image_base64 = encode_image_to_base64(image)
        
        client = anthropic.Anthropic(api_key=api_key, timeout=60.0)
        
        def make_call():
            return client.messages.create(
                model=model,
                max_tokens=8192,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "image",
                                "source": {
                                    "type": "base64",
                                    "media_type": "image/png",
                                    "data": image_base64
                                }
                            },
                            {"type": "text", "text": prompt}
                        ]
                    }
                ]
            )
        
        message = await asyncio.get_event_loop().run_in_executor(executor, make_call)
        response_text = message.content[0].text
        
        # Capture usage statistics from Claude API
        usage = {
            "input_tokens": message.usage.input_tokens,
            "output_tokens": message.usage.output_tokens,
            "total_tokens": message.usage.input_tokens + message.usage.output_tokens
        }
        
        # Extract JSON
        json_text = response_text.strip()
        if "```json" in json_text:
            start = json_text.find("```json") + 7
            end = json_text.find("```", start)
            json_text = json_text[start:end].strip()
        elif "```" in json_text:
            start = json_text.find("```") + 3
            end = json_text.find("```", start)
            json_text = json_text[start:end].strip()
        
        try:
            result = json.loads(json_text)
            parsed = _parse_vision_response(result)
            if parsed:
                parsed["usage"] = usage  # Add usage data
            return parsed

        except json.JSONDecodeError as e:
            print(f"⚠️ JSON parse error (likely truncated): {e}")
            # First attempt: repair the truncated JSON
            repaired = _repair_truncated_json(json_text)
            if repaired:
                try:
                    result = json.loads(repaired)
                    print(f"   ✅ JSON repaired successfully")
                    return _parse_vision_response(result)
                except json.JSONDecodeError:
                    print(f"   ⚠️ Repair failed, trying API retry...")
            # Retry with a compact prompt to reduce output size
            print(f"   🔄 Retrying with compact prompt...")
            try:
                compact_prompt = _create_compact_vision_prompt(source_lang, target_lang)
                def make_compact_call():
                    return client.messages.create(
                        model=model,
                        max_tokens=8192,
                        messages=[{
                            "role": "user",
                            "content": [
                                {"type": "image", "source": {"type": "base64", "media_type": "image/png", "data": image_base64}},
                                {"type": "text", "text": compact_prompt}
                            ]
                        }]
                    )
                retry_message = await asyncio.get_event_loop().run_in_executor(executor, make_compact_call)
                retry_text = retry_message.content[0].text.strip()
                if "```json" in retry_text:
                    start = retry_text.find("```json") + 7
                    end = retry_text.find("```", start)
                    retry_text = retry_text[start:end].strip()
                elif "```" in retry_text:
                    start = retry_text.find("```") + 3
                    end = retry_text.find("```", start)
                    retry_text = retry_text[start:end].strip()
                result = json.loads(retry_text)
                print(f"   ✅ Compact retry succeeded")
                return _parse_vision_response(result)
            except Exception as retry_err:
                print(f"   ❌ Compact retry also failed: {retry_err}")
                return None
    
    except Exception as e:
        error_msg = str(e)
        # Classify known API errors for user-friendly messages
        if "credit balance" in error_msg or "402" in error_msg:
            user_msg = "Insufficient credits on your Anthropic account. Please top up at console.anthropic.com."
        elif "invalid_api_key" in error_msg or "401" in error_msg or "authentication" in error_msg.lower():
            user_msg = "Invalid Claude API key. Please check your key in Settings."
        elif "Connection error" in error_msg or "ConnectError" in error_msg:
            user_msg = "Cannot reach Anthropic API. Check your internet connection."
        elif "overloaded" in error_msg or "529" in error_msg:
            user_msg = "Anthropic API is temporarily overloaded. Please retry in a few moments."
        elif "rate_limit" in error_msg or "429" in error_msg:
            user_msg = "Anthropic API rate limit reached. Please wait before retrying."
        else:
            user_msg = f"Claude API error: {error_msg}"
        print(f"❌ {user_msg}")
        import traceback
        traceback.print_exc()
        return {"error": user_msg}


async def call_openrouter_vision(
    image: Image.Image,
    prompt: str,
    api_key: Optional[str],
    model: str = "google/gemini-flash-1.5"
) -> Optional[Dict]:
    """Call OpenRouter API with vision"""
    try:
        image_base64 = encode_image_to_base64(image)

        headers = {"Content-Type": "application/json"}
        if api_key:
            headers["Authorization"] = f"Bearer {api_key}"

        payload = {
            "model": model,
            "messages": [{
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_base64}"}},
                    {"type": "text", "text": prompt}
                ]
            }]
        }

        max_retries = 3
        backoff = 2
        response = None
        for attempt in range(max_retries + 1):
            response = await asyncio.get_event_loop().run_in_executor(
                executor,
                lambda: requests.post(
                    "https://openrouter.ai/api/v1/chat/completions",
                    headers=headers,
                    json=payload,
                    timeout=60
                )
            )
            if response.status_code == 429 and attempt < max_retries:
                wait_s = backoff ** (attempt + 1)
                print(f"⚠️ OpenRouter rate limit (429). Retrying in {wait_s}s...")
                await asyncio.sleep(wait_s)
                continue
            response.raise_for_status()
            break

        if response is None:
            return None

        result = response.json()
        
        # Capture usage statistics from API response
        usage = result.get("usage", {})
        input_tokens = usage.get("input_tokens", 0)
        output_tokens = usage.get("output_tokens", 0)
        
        response_text = result["choices"][0]["message"]["content"]

        if "```json" in response_text:
            start = response_text.find("```json") + 7
            end = response_text.find("```", start)
            json_text = response_text[start:end].strip()
        else:
            json_text = response_text.strip()

        result = json.loads(json_text)
        
        parsed_result = _parse_vision_response(result)
        
        # Add usage data to result
        if parsed_result:
            parsed_result["usage"] = {
                "input_tokens": input_tokens,
                "output_tokens": output_tokens,
                "total_tokens": input_tokens + output_tokens
            }

        return parsed_result
    
    except requests.exceptions.Timeout:
        print(f"⏱️ OpenRouter vision timeout after 60s")
        return {"error": "API timeout - the translation service took too long to respond. Please try again."}
    except requests.exceptions.ConnectionError:
        print(f"🔌 OpenRouter vision connection error")
        return {"error": "Connection error - unable to reach the translation service. Please check your internet connection."}
    except Exception as e:
        print(f"❌ OpenRouter vision error: {e}")
        return {"error": f"Translation error: {str(e)}"}



def ocr_extract_text(image: Image.Image, source_lang: str = "es") -> str:
    """Extract raw text from image using pytesseract (local, free).
    PSM 3 = fully automatic page segmentation (better for complex slide layouts).
    Language code passed to tesseract for better accuracy.
    """
    try:
        # Map our lang codes to tesseract lang codes
        tess_lang = {"es": "spa", "fr": "fra", "it": "ita", "en": "eng"}.get(source_lang, "spa+eng")

        # Upscale small images for better OCR accuracy (tesseract works best at 300dpi+)
        w, h = image.size
        if w < 1200:
            scale = 1200 / w
            image = image.resize((int(w * scale), int(h * scale)), Image.LANCZOS)

        # Convert to grayscale + light contrast boost
        import PIL.ImageEnhance
        image = image.convert("L")
        image = PIL.ImageEnhance.Contrast(image).enhance(1.5)

        # PSM 3: fully automatic page segmentation, no OSD
        # Better than PSM 6 for slides with mixed columns/layouts
        config = f"--psm 3 --oem 3 -l {tess_lang}"
        text = pytesseract.image_to_string(image, config=config)

        # Clean: collapse multiple blank lines into one, strip trailing spaces per line
        lines = text.splitlines()
        cleaned_lines = []
        prev_empty = False
        for line in lines:
            stripped = line.strip()
            if stripped:
                cleaned_lines.append(stripped)
                prev_empty = False
            elif not prev_empty:
                cleaned_lines.append("")
                prev_empty = True

        raw = "\n".join(cleaned_lines).strip()

        # Remove emoji and non-printable characters that OCR misreads
        import re
        raw = re.sub(
            r'[\U00010000-\U0010FFFF'   # Supplementary planes (most emojis)
            r'\U00002600-\U000027BF'    # Misc symbols (☆, ©, ✓, etc.)
            r'\U0001F300-\U0001F9FF'    # Emoji blocks
            r']', '', raw
        )
        
        # Clean up lines: remove garbage and scale artifacts
        def is_scale_artifact(text):
            """Detect visual scale artifacts misread by OCR."""
            # Pattern 1: Sequence of numbers with spaces (may have trailing ], ), etc.)
            if re.search(r'(\d+\s+){3,}\d*[\]\)]?', text):
                alphas = len(re.findall(r'[a-zA-Z]', text))
                if alphas <= 2:  # Max 2 letters allowed
                    return True
            
            # Pattern 2: Sequence of short words (1-3 letters, uppercase)
            words = text.split()
            if len(words) >= 5:
                short_words = sum(1 for w in words if len(w) <= 3 and w.isupper())
                if short_words >= len(words) * 0.8:  # 80%+ are short uppercase
                    return True
            
            return False
        
        def is_garbage_line(text):
            """Detect OCR garbage lines that should be removed."""
            # Very short line (<15 chars) with lots of special chars or mixed case nonsense
            if len(text) < 15:
                # Count letters vs special chars
                letters = len(re.findall(r'[a-zA-Z]', text))
                special = len(re.findall(r'[^a-zA-Z0-9\s]', text))
                # If more special chars than letters, it's garbage
                if special > letters:
                    return True
                # Common garbage patterns
                if re.match(r'^[a-z]{1,3}\.\s*$', text):  # "lnire.", "x f."
                    return True
            return False
        
        def is_bullet(text):
            """Detect if line is a bullet point."""
            return bool(re.match(r'^[+#\-•*◦▪▫]\s+|^[a-z]\s+', text))
        
        def is_title(text):
            """Detect if line is likely a title."""
            # Starts with number pattern
            if re.match(r'^\d+\s*-\s*', text):
                return True
            # Contains title keywords
            if any(kw in text for kw in ['(Option)', '(Scale)', '(Question)', 'Intro']):
                return True
            return False
        
        # First pass: filter out garbage and artifacts
        cleaned = []
        for ln in raw.splitlines():
            stripped = ln.strip()
            
            if not stripped:
                cleaned.append("")
                continue
            
            if not re.search(r'[a-zA-ZÀ-ÿ0-9]', stripped):
                continue
            
            if is_scale_artifact(stripped):
                continue
            
            if is_garbage_line(stripped):
                continue
            
            cleaned.append(stripped)
        
        # Second pass: smart grouping
        lines_out = []
        i = 0
        while i < len(cleaned):
            line = cleaned[i]
            
            if not line:
                lines_out.append("")
                i += 1
                continue
            
            # Title → always separate with blank line before/after
            if is_title(line):
                if lines_out and lines_out[-1] != "":
                    lines_out.append("")  # Blank before title
                lines_out.append(line)
                i += 1
                continue
            
            # Bullet → keep separate, no blank lines between bullets
            if is_bullet(line):
                lines_out.append(line)
                i += 1
                continue
            
            # Short line after question (likely radio option) → keep separate
            if i > 0 and cleaned[i-1].endswith('?') and len(line) < 80:
                lines_out.append(line)
                i += 1
                continue
            
            # Regular line → try to join with next lines if they're continuations
            paragraph = [line]
            j = i + 1
            while j < len(cleaned):
                next_line = cleaned[j]
                
                # Stop at empty line
                if not next_line:
                    break
                
                # Stop at title or bullet
                if is_title(next_line) or is_bullet(next_line):
                    break
                
                current_text = " ".join(paragraph)
                
                # SIMPLE RULE: If current text doesn't end with strong punctuation, ALWAYS join
                # (unless next line is clearly a separate item)
                if not current_text.rstrip().endswith(('.', '!', '?', ':')):
                    # Exception: if next line starts with quote and is short, might be radio option
                    if next_line.startswith('"') and len(next_line) < 80:
                        # Check if there are more short quoted lines (indicating list)
                        looks_like_list = False
                        for k in range(j, min(j+3, len(cleaned))):
                            if cleaned[k] and (cleaned[k].startswith('"') or 
                               (len(cleaned[k]) < 60 and k == j)):
                                looks_like_list = True
                        if looks_like_list:
                            break
                    # Otherwise, join
                    paragraph.append(next_line)
                    j += 1
                    continue
                
                # Current text ends with punctuation
                # Join if next line is clearly a continuation (lowercase or punctuation start)
                if next_line[0].islower() or next_line[0] in '",;:':
                    paragraph.append(next_line)
                    j += 1
                    continue
                
                # Check if next line looks like start of a new paragraph (long line with uppercase)
                if next_line[0].isupper() and len(next_line) > 60:
                    break  # New paragraph
                
                # Check if next line looks like start of a list (short lines following)
                if next_line[0].isupper() and len(next_line) < 60:
                    looks_like_list = False
                    for k in range(j, min(j+3, len(cleaned))):
                        if cleaned[k] and len(cleaned[k]) < 60:
                            looks_like_list = True
                    if looks_like_list:
                        break  # It's a list
                
                # Otherwise join
                paragraph.append(next_line)
                j += 1
            
            # Output joined paragraph
            lines_out.append(" ".join(paragraph))
            i = j
        
        lines_out = []
        for ln in raw.splitlines():
            stripped = ln.strip()
            
            # Skip empty lines (will be preserved as-is)
            if not stripped:
                lines_out.append("")
                continue
            
            # Skip pure garbage lines (no letters/digits, only symbols/spaces)
            if not re.search(r'[a-zA-ZÀ-ÿ0-9]', stripped):
                continue
            
            # Skip visual scale artifacts
            if is_scale_artifact(stripped):
                continue
            
            # Skip garbage lines
            if is_garbage_line(stripped):
                continue
            
            lines_out.append(ln)
        
        return "\n".join(lines_out).strip()

    except Exception as e:
        print(f"❌ OCR error: {e}")
        return ""


def detect_language_from_image(image: Image.Image) -> Optional[str]:
    """Detect language from an image using local OCR + langdetect."""
    text = ocr_extract_text(image, source_lang="es")
    if not text or len(text.strip()) < 20:
        return None
    try:
        lang = detect(text)
    except LangDetectException:
        return None
    if lang in {"es", "fr", "it", "en"}:
        return lang
    return None


def detect_language_from_presentation(prs, selected_slides: List[int]) -> Optional[str]:
    """Detect language from the first selected slide that has an image."""
    selected_set = set(selected_slides)
    for idx, slide in enumerate(prs.slides):
        if idx not in selected_set:
            continue
        has_img, image = has_image_on_left(slide)
        if not has_img or image is None:
            continue
        lang = detect_language_from_image(image)
        if lang:
            return lang
    return None


# Cache of loaded CTranslate2 translators (avoid reloading per slide)
_ct2_translators: Dict[str, any] = {}


def _get_ct2_translator(source_lang: str, target_lang: str):
    """
    Load (or return cached) CTranslate2 Helsinki-NLP OPUS-MT translator.
    Uses ct2-transformers-converter (CLI) — the correct approach for HuggingFace
    format models. OpusMTConverter only works with old Marian .yml format.
    Downloads and converts on first use (~300MB per pair), cached after that.
    """
    import os
    import subprocess
    import tempfile
    import ctranslate2
    import transformers

    pair_key = f"{source_lang}-{target_lang}"
    if pair_key in _ct2_translators:
        return _ct2_translators[pair_key]

    hf_model_name = f"Helsinki-NLP/opus-mt-{source_lang}-{target_lang}"
    ct2_dir = os.path.join(tempfile.gettempdir(), f"ct2_opus_{source_lang}_{target_lang}")

    if not os.path.exists(os.path.join(ct2_dir, "model.bin")):
        print(f"   📦 Downloading + converting {hf_model_name} (one-time ~300MB)...")
        result = subprocess.run(
            [
                "ct2-transformers-converter",
                "--model", hf_model_name,
                "--output_dir", ct2_dir,
                "--quantization", "int8",
                "--force",
                "--copy_files", "source.spm", "target.spm",
                "vocab.json", "tokenizer_config.json", "generation_config.json"
            ],
            capture_output=True, text=True
        )
        if result.returncode != 0:
            raise RuntimeError(
                f"ct2-transformers-converter failed:\n{result.stderr}\n"
                f"Make sure it is installed: pip install ctranslate2"
            )
        print(f"   ✅ Conversion done → {ct2_dir}")
    else:
        print(f"   ✅ Loading cached model from {ct2_dir}")

    translator = ctranslate2.Translator(ct2_dir, device="cpu", inter_threads=2)
    # Load tokenizer from HuggingFace (not from ct2_dir — requires sentencepiece)
    tokenizer = transformers.AutoTokenizer.from_pretrained(hf_model_name)

    _ct2_translators[pair_key] = (translator, tokenizer)
    print(f"   ✅ Model ready and cached in memory")
    return translator, tokenizer


def translate_offline(text: str, source_lang: str, target_lang: str) -> Optional[str]:
    """
    Offline translation using CTranslate2 + Helsinki-NLP OPUS-MT models.
    - Quality: significantly better than ArgosTranslate
    - Speed: fast (C++ int8, no GPU needed)
    - Size: ~300MB per language pair, downloaded once
    - Supported: es/fr/it → en
    """
    try:
        translator, tokenizer = _get_ct2_translator(source_lang, target_lang)

        lines = text.splitlines()
        translated_lines = []

        for line in lines:
            if not line.strip():
                translated_lines.append("")
                continue
            # Tokenize with HuggingFace tokenizer (handles SentencePiece internally)
            tokens = tokenizer.convert_ids_to_tokens(tokenizer.encode(line))
            results = translator.translate_batch([tokens])
            translated_ids = tokenizer.convert_tokens_to_ids(results[0].hypotheses[0])
            translated = tokenizer.decode(translated_ids, skip_special_tokens=True)
            translated_lines.append(translated)

        return "\n".join(translated_lines)

    except ImportError as e:
        print(f"   ❌ Missing dependency: {e}")
        print(f"   → Run: pip install ctranslate2 transformers sentencepiece")
        return None
    except Exception as e:
        print(f"   ❌ Offline translation error: {e}")
        import traceback
        traceback.print_exc()
        return None


def build_text_structure(translated_text: str) -> List[Dict]:
    """Convert translated text into paragraph structure for apply_formatting_from_structure.
    Each non-empty line becomes a paragraph. Empty lines become spacers.
    Collapses consecutive empty lines into one.
    """
    structure = []
    prev_empty = False

    for line in translated_text.splitlines():
        stripped = line.strip()

        # Collapse consecutive empty lines
        if not stripped:
            if prev_empty:
                continue
            prev_empty = True
        else:
            prev_empty = False

        structure.append({
            "type": "paragraph",
            "original": stripped,
            "translated": stripped,
            "formatting": {
                "bold": False,
                "size_relative": 0,
                "color": "black",
                "alignment": "left"
            },
            "has_emoji": False
        })

    return structure



def post_process_translation(text: str) -> str:
    """
    Clean up LLM translation output to fix structure issues.
    The free LLM models often ignore prompt instructions, so we enforce them here.
    """
    import re
    
    lines = text.splitlines()
    result = []
    i = 0
    
    while i < len(lines):
        line = lines[i].strip()
        
        # Empty line - preserve
        if not line:
            result.append("")
            i += 1
            continue
        
        # Check if this line should be joined with next line
        if i + 1 < len(lines):
            next_line = lines[i + 1].strip()
            
            if next_line:
                # Don't join if next line is clearly a new section
                is_next_title = bool(re.match(r'^\d+\s*-', next_line))
                is_next_bullet = bool(re.match(r'^[+#\-•*◦▪▫$]\s+', next_line))
                is_next_option = next_line.startswith('"') and len(next_line) < 80
                
                # Don't join if next line contains title keywords
                has_title_keyword = any(kw in next_line for kw in ['(Option)', '(Scale)', '(Question)', 'Intro'])
                
                # Don't join if current line ends with a name/signature pattern
                # (name, title/company format)
                is_signature = bool(re.search(r',\s+[A-Z][a-z]+.*(?:contact|representative|team|support)$', line, re.IGNORECASE))
                
                # Don't join if next line starts with uppercase AND current line is reasonably complete
                # (>40 chars AND has at least one verb/subject structure)
                looks_complete = len(line) > 40 and line.count(' ') > 5
                
                should_not_join = (
                    is_next_title or 
                    is_next_bullet or 
                    is_next_option or 
                    has_title_keyword or
                    is_signature
                )
                
                # JOIN if current line clearly continues (doesn't end with punctuation)
                # AND next line is not a new section
                if not line.endswith(('.', '!', '?', ':')) and not should_not_join:
                    # Additional check: if next line starts with lowercase, definitely join
                    if next_line[0].islower() or next_line[0] in ',;':
                        result.append(line + " " + next_line)
                        i += 2
                        continue
                    
                    # If next line starts with uppercase but is long (>60 chars), it might be a new paragraph
                    # But if current line is short (<60 chars), it's probably incomplete
                    if len(line) < 60 or len(next_line) < 60:
                        result.append(line + " " + next_line)
                        i += 2
                        continue
        
        result.append(line)
        i += 1
    
    # Second pass: remove blank lines in lists
    final = []
    i = 0
    while i < len(result):
        line = result[i]
        
        # If this is a blank line, check if we're in a list
        if not line:
            # Look at previous and next non-empty lines
            prev_line = None
            next_line = None
            
            for j in range(i-1, -1, -1):
                if result[j]:
                    prev_line = result[j]
                    break
            
            for j in range(i+1, len(result)):
                if result[j]:
                    next_line = result[j]
                    break
            
            # Check if both prev and next are list items
            if prev_line and next_line:
                is_prev_bullet = bool(re.match(r'^[+#\-•*◦▪▫$]\s+', prev_line))
                is_next_bullet = bool(re.match(r'^[+#\-•*◦▪▫$]\s+', next_line))
                
                is_prev_option = (prev_line.startswith('"') or 
                                 (len(prev_line) < 80 and not prev_line.endswith(('?', '.', '!'))))
                is_next_option = (next_line.startswith('"') or 
                                 (len(next_line) < 80 and not next_line.endswith(('?', '.', '!'))))
                
                # If both are bullets or both are options, skip the blank line
                if (is_prev_bullet and is_next_bullet) or (is_prev_option and is_next_option):
                    i += 1
                    continue
        
        final.append(line)
        i += 1
    
    return "\n".join(final)


async def call_ocr_openrouter(
    image: Image.Image,
    source_lang: str,
    target_lang: str,
    api_key: str,
    model: str = "mistralai/mistral-7b-instruct:free",
    log_file: Optional[str] = None
) -> Optional[Dict]:
    """
    Free pipeline: local pytesseract OCR → OpenRouter text LLM translation.
    Falls back to ArgosTranslate (offline) if OpenRouter fails.
    Returns translation_method: "openrouter" | "offline" | None
    No formatting preserved — plain black text output.
    """
    # Log helper — writes to terminal AND optional file
    def log(msg: str):
        print(msg)
        if log_file:
            with open(log_file, "a", encoding="utf-8") as f:
                f.write(msg + "\n")

    # Step 1: OCR locally (always free)
    log(f"   🔍 OCR extracting text (lang={source_lang})...")
    raw_text = ocr_extract_text(image, source_lang)

    if not raw_text:
        log(f"   ⚠️ OCR returned empty text")
        return None

    log(f"   ✅ OCR extracted {len(raw_text)} chars")

    lang_names = {"fr": "French", "es": "Spanish", "it": "Italian", "en": "English"}
    src = lang_names.get(source_lang, source_lang)
    tgt = lang_names.get(target_lang, target_lang)

    # Step 2a: Try OpenRouter first
    openrouter_error = None
    translated_text = None
    translation_method = None
    translation_model = None

    if api_key:
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}"
        }
        system_prompt = (
            f"You are a translator. Translate from {src} to {tgt}. "
            "\n"
            "The text has been pre-processed, but may still have formatting issues. Your job:\n"
            "\n"
            "1. TRANSLATE accurately\n"
            "\n"
            "2. FIX remaining structure issues:\n"
            "   - If you see a paragraph broken across multiple lines (incomplete sentences), JOIN them\n"
            "   - If you see a list (bullets or options), ensure NO blank lines between items\n"
            "   - Keep blank lines ONLY between distinct sections/paragraphs\n"
            "\n"
            "3. DETECT lists:\n"
            "   - Bullet lists (lines starting with #, +, -, *, e, etc.) → no blank lines between items\n"
            "   - Radio options (short lines after a question) → no blank lines between items\n"
            "\n"
            "Examples:\n"
            "BAD: 'Summer is joy\\nand travel.' → GOOD: 'Summer is joy and travel.'\n"
            "BAD: '# Item 1\\n\\n# Item 2' → GOOD: '# Item 1\\n# Item 2'\n"
            "\n"
            "Output ONLY the translated, properly structured text."
        )

        # Some providers/models (e.g., Gemma via Google AI Studio) do not allow system messages.
        use_system = "gemma" not in model.lower()
        if use_system:
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": raw_text}
            ]
        else:
            messages = [
                {"role": "user", "content": system_prompt + "\n\nText to translate:\n" + raw_text}
            ]

        payload = {
            "model": model,
            "messages": messages,
            "max_tokens": 3000,
            "temperature": 0.1
        }

        log(f"   🌐 Translating with {model}...")
        max_retries = 3
        backoff = 2
        response = None

        try:
            for attempt in range(max_retries + 1):
                response = await asyncio.get_event_loop().run_in_executor(
                    executor,
                    lambda: requests.post(
                        "https://openrouter.ai/api/v1/chat/completions",
                        headers=headers,
                        json=payload,
                        timeout=60
                    )
                )
                if response.status_code == 429 and attempt < max_retries:
                    wait_s = backoff ** (attempt + 1)
                    log(f"   ⚠️ Rate limit (429). Retrying in {wait_s}s... ({attempt + 1}/{max_retries})")
                    await asyncio.sleep(wait_s)
                    continue
                if response.status_code >= 400:
                    preview = response.text[:800] if response.text else ""
                    log(f"   ❌ OpenRouter HTTP {response.status_code}: {preview}")
                response.raise_for_status()
                break

            if response and response.ok:
                result = response.json()
                raw_response = result["choices"][0]["message"]["content"].strip()
                if raw_response:
                    # Post-process the LLM output to fix structure issues it ignored
                    translated_text = post_process_translation(raw_response)
                    translation_method = "openrouter"
                    translation_model = model
                    log(f"   ✅ OpenRouter translation OK: {len(translated_text)} chars")
                else:
                    openrouter_error = "Empty response from model"
                    log(f"   ⚠️ OpenRouter returned empty response")

        except Exception as e:
            # If we have a response body, surface it for easier debugging
            if response is not None and getattr(response, "text", None):
                openrouter_error = response.text
            else:
                openrouter_error = str(e)
            log(f"   ⚠️ OpenRouter failed: {e}")
    else:
        openrouter_error = "No API key configured"

    # Step 2b: Fallback to offline CTranslate2 + Helsinki-NLP
    if not translated_text:
        # If source == target, no translation needed — return raw text directly
        if source_lang == target_lang:
            log(f"   ℹ️ Source == target language ({source_lang}), skipping translation")
            translated_text = raw_text
            translation_method = "passthrough"
            translation_model = "passthrough"
        else:
            log(f"   💻 Falling back to offline translation (Helsinki-NLP)...")
            translated_text = await asyncio.get_event_loop().run_in_executor(
                executor,
                lambda: translate_offline(raw_text, source_lang, target_lang)
            )
            if translated_text:
                translation_method = "offline"
                translation_model = f"Helsinki-NLP/opus-mt-{source_lang}-{target_lang}"
                log(f"   ✅ Offline translation OK: {len(translated_text)} chars")
            else:
                log(f"   ❌ Offline translation also failed")
                return None

    # Step 3: Build structure
    # Guard: if translation produced only empty lines, treat as failure
    if not translated_text or not translated_text.strip():
        log(f"   ❌ Translation result is empty after parsing — skipping slide")
        return None

    structure = build_text_structure(translated_text)

    return {
        "original_text": raw_text,
        "translated_text": translated_text,
        "structure": structure,
        "translation_method": translation_method,           # "openrouter" or "offline"
        "translation_model": translation_model,
        "openrouter_error": openrouter_error,               # None if openrouter succeeded
    }


def has_image_on_left(slide) -> Tuple[bool, Optional[Image.Image]]:
    """Check if slide has image on left side"""
    try:
        if hasattr(slide, 'part') and hasattr(slide.part, 'presentation_part'):
            slide_width = slide.part.presentation_part.presentation.slide_width
        else:
            slide_width = 9144000
    except:
        slide_width = 9144000
    
    for shape in slide.shapes:
        if shape.shape_type == 13:  # Picture
            if shape.left < slide_width * 0.5:
                try:
                    image_stream = io.BytesIO(shape.image.blob)
                    image = Image.open(image_stream)
                    return True, image
                except Exception as e:
                    print(f"Image error: {e}")
    
    return False, None


def reduce_font_if_overflow(slide, max_iterations=20, min_font_size=6):
    """
    Manually reduce font size for all text in a slide until content fits within the textbox.
    PowerPoint's TEXT_TO_FIT_SHAPE doesn't always reduce enough, so we do it manually.
    
    We can't directly check if text overflows, but we can compare the text height
    with the frame height and keep reducing until it fits.
    """
    from pptx.util import Pt
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        text_frame = shape.text_frame
        
        # Skip if no text
        if not any(p.text.strip() for p in text_frame.paragraphs):
            continue
        
        # Try to access the actual text height (this is approximate)
        # We'll use a heuristic: keep reducing while we have "too much" text
        for iteration in range(max_iterations):
            reduced = False
            smallest_font = 100  # Track smallest font in use
            
            # Collect all current font sizes
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        smallest_font = min(smallest_font, run.font.size.pt)
            
            # If we've hit the minimum, stop
            if smallest_font <= min_font_size:
                break
            
            # Estimate if we need to reduce: count lines vs available height
            num_paragraphs = len([p for p in text_frame.paragraphs if p.text.strip()])
            total_chars = sum(len(p.text) for p in text_frame.paragraphs)
            
            # Rough heuristic: if we have >30 paragraphs or >2500 chars, we likely overflow
            # This threshold can be adjusted based on testing
            avg_chars_per_para = total_chars / max(num_paragraphs, 1)
            needs_reduction = (
                num_paragraphs > 30 or 
                total_chars > 2500 or
                (num_paragraphs > 20 and avg_chars_per_para > 60)
            )
            
            if not needs_reduction:
                break
            
            # Reduce all fonts by 0.5pt
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size and run.font.size.pt > min_font_size:
                        current_size = run.font.size.pt
                        new_size = max(min_font_size, current_size - 0.5)
                        run.font.size = Pt(new_size)
                        reduced = True
            
            # If nothing was reduced (all at minimum), break
            if not reduced:
                break
        
        # After all reductions, ensure word wrap is on and margins are minimal
        text_frame.word_wrap = True
        from pptx.util import Inches
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        text_frame.margin_left = Inches(0.08)
        text_frame.margin_right = Inches(0.08)


def apply_formatting_from_structure(
    slide,
    structure: List[Dict],
    left: float,
    top: float,
    width: float,
    height: float,
    base_font_size: int,
    title_adjustment: int,
    subject_adjustment: int,
    ocr_mode: bool = False
):
    """Apply formatted text to slide with proper spacing"""
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE  # Reduce text size if it overflows
    text_frame.margin_top = Inches(0.1)
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_bottom = Inches(0.1)
    
    color_map = {
        "green": RGBColor(46, 125, 50),
        "blue": RGBColor(25, 118, 210),
        "red": RGBColor(211, 47, 47),
        "orange": RGBColor(245, 124, 0),
        "black": RGBColor(33, 33, 33),
        "grey": RGBColor(117, 117, 117),
        "gray": RGBColor(117, 117, 117),
    }

    # Line-based rendering (Claude vision mode)
    if structure and isinstance(structure[0], dict) and "line_blocks" in structure[0]:

        def set_para_style(para, alignment):
            if alignment == "center":
                para.alignment = PP_ALIGN.CENTER
            elif alignment == "right":
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.LEFT
            para.line_spacing = 1.15

        def write_run(para, text, formatting):
            clean = text.rstrip("\n")
            if not clean:
                return
            run = para.add_run()
            run.text = clean
            is_bold = formatting.get("bold", False)
            size_relative = formatting.get("size_relative", 0)
            run.font.bold = is_bold
            run.font.size = Pt(base_font_size + size_relative)
            run.font.name = "Calibri"
            color_name = formatting.get("color", "black").lower()
            if color_name == "red" and not is_bold:
                color_name = "black"
            run.font.color.rgb = color_map.get(color_name, RGBColor(33, 33, 33))

        para_idx = 0
        for line in structure:
            alignment = line.get("alignment", "left")
            line_blocks = line.get("line_blocks", [])

            # Empty line → spacer paragraph (skip consecutive spacers)
            if not line_blocks:
                last_para = text_frame.paragraphs[-1] if text_frame.paragraphs else None
                last_is_spacer = last_para and all(r.text.strip() in ("", " ") for r in last_para.runs)
                if last_is_spacer and para_idx > 0:
                    continue
                p = text_frame.paragraphs[0] if para_idx == 0 else text_frame.add_paragraph()
                para_idx += 1
                set_para_style(p, alignment)
                run = p.add_run()
                run.text = " "
                run.font.size = Pt(base_font_size)
                run.font.name = "Calibri"
                continue

            # Content line: start one paragraph, but each block ending with \n
            # starts a fresh paragraph for the next block
            p = text_frame.paragraphs[0] if para_idx == 0 else text_frame.add_paragraph()
            para_idx += 1
            set_para_style(p, alignment)

            # Filter out empty blocks, then iterate with index to know if last
            active_blocks = [b for b in line_blocks if b.get("translated", "")]
            for b_idx, block in enumerate(active_blocks):
                text = block.get("translated", "")
                formatting = block.get("formatting", {})
                is_last_block = (b_idx == len(active_blocks) - 1)
                ends_with_newline = text.endswith("\n")
                # Add space between adjacent blocks on the same paragraph
                # (e.g. scale numbers "1","2","3" should read "1 2 3")
                if p.runs and not p.runs[-1].text.endswith((" ", "\n")) and not text.lstrip("\n").startswith((" ", ".", ",", "!", "?", ":", ";", ")", "]")):
                    text = " " + text
                write_run(p, text, formatting)
                # Only break to new paragraph on \n if there are more blocks after this one
                if ends_with_newline and not is_last_block:
                    p = text_frame.add_paragraph()
                    para_idx += 1
                    set_para_style(p, alignment)

        return

    # OCR free mode: each block = one paragraph (one line of text)
    if ocr_mode:
        for idx, block in enumerate(structure):
            text = block.get("translated", "")
            formatting = block.get("formatting", {})

            para = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()

            alignment = formatting.get("alignment", "left")
            if alignment == "center":
                para.alignment = PP_ALIGN.CENTER
            elif alignment == "right":
                para.alignment = PP_ALIGN.RIGHT
            else:
                para.alignment = PP_ALIGN.LEFT
            para.line_spacing = 1.15

            if not text or not text.strip():
                run = para.add_run()
                run.text = ""
                run.font.size = Pt(base_font_size)
                run.font.name = "Calibri"
                continue

            run = para.add_run()
            run.text = text.strip()
            is_bold = formatting.get("bold", False)
            size_relative = formatting.get("size_relative", 0)
            run.font.bold = is_bold
            run.font.size = Pt(base_font_size + size_relative)
            run.font.name = "Calibri"
            color_name = formatting.get("color", "black").lower()
            if color_name == "red" and not is_bold:
                color_name = "black"
            run.font.color.rgb = color_map.get(color_name, RGBColor(33, 33, 33))
        return

async def process_ppt_slide(
    slide,
    source_lang: str,
    target_lang: str,
    provider: str,
    api_key: Optional[str],
    claude_api_key: Optional[str],
    model: str,
    base_font_size: int,
    title_adjustment: int,
    subject_adjustment: int,
    log_file: Optional[str] = None,
    slide_num: int = 0
) -> Optional[Dict]:
    """Process single slide"""
    
    if log_file:
        with open(log_file, "a", encoding="utf-8") as f:
            f.write(f"\n{'='*60}\nSLIDE {slide_num}\n{'='*60}\n")
    print(f"🔍 Processing slide with {provider}")
    
    has_img, image = has_image_on_left(slide)
    if not has_img or image is None:
        print(f"⚠️ No image on left")
        return None
        
    if provider == "claude":
        prompt = create_vision_prompt(source_lang, target_lang)
        result = await call_claude_vision(image, prompt, api_key, model)
    elif provider == "openrouter":
        prompt = create_vision_prompt(source_lang, target_lang)
        result = await call_openrouter_vision(image, prompt, api_key, model)
        if result is None:
            print("⚠️ OpenRouter failed, falling back to OCR + offline translation...")
            # Force offline translation by skipping OpenRouter in the OCR pipeline
            result = await call_ocr_openrouter(image, source_lang, target_lang, api_key=None, model=model, log_file=log_file)
    elif provider == "ocr_free":
        # Free pipeline: local OCR + OpenRouter text LLM (no vision needed)
        result = await call_ocr_openrouter(image, source_lang, target_lang, api_key, model, log_file=log_file)
    else:
        print(f"❌ Unknown provider: {provider}")
        return None
    
    if result and isinstance(result, dict) and "error" in result and "structure" not in result:
        print(f"❌ API error: {result['error']}")
        return result  # propagate error dict upstream
    if result:
        print(f"✅ Translation received")
    else:
        print(f"❌ No result from API")
    
    return result


async def create_translated_ppt(
    input_file_bytes: bytes,
    selected_slides: List[int],
    source_lang: str,
    target_lang: str,
    provider: str,
    api_key: Optional[str],
    model: str,
    settings: Dict,
    progress_callback=None
) -> Tuple[bytes, Dict]:
    """Create translated PowerPoint"""
    
    import datetime
    import os
    # Create a log file next to the output, timestamped
    log_timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    log_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        f"ocr_debug_{log_timestamp}.txt"
    )
    with open(log_path, "w", encoding="utf-8") as f:
        f.write(f"OCR + Translation debug log — {log_timestamp}\n")
        f.write(f"Source: {source_lang} → Target: {target_lang} | Model: {model}\n")
        f.write("=" * 60 + "\n")
    print(f"📄 Debug log: {log_path}")

    prs = Presentation(io.BytesIO(input_file_bytes))
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height
    
    # Number of selected slides for progress tracking
    selected_slides_count = len(selected_slides)
    
    stats = {
        "total_slides": len(prs.slides),
        "selected_slides_count": selected_slides_count,  # For accurate progress
        "processed_slides": 0,
        "skipped_slides": 0,
        "failed_slides": 0,
        "slide_methods": [],  # List of {slide: int, method: str, error: str|None}
        "warnings": [],       # User-facing warning messages
        "total_input_tokens": 0,
        "total_output_tokens": 0,
        "total_cost": 0.0,
        "start_time": time.time()  # Track start time
    }
    
    base_font_size = settings.get("base_font_size", 11)
    title_adjustment = settings.get("title_size_adjustment", 5)
    subject_adjustment = settings.get("subject_size_adjustment", 1)
    claude_api_key = settings.get("claude_api_key")
    
    # Counter for selected slides processed (1, 2, 3...) instead of slide numbers (15, 16, 23...)
    selected_slide_counter = 0
    
    for idx, slide in enumerate(prs.slides):
        if idx not in selected_slides:
            stats["skipped_slides"] += 1
            # Don't report progress for skipped slides
            continue
        
        # Increment counter for selected slides
        selected_slide_counter += 1
        
        if progress_callback:
            await progress_callback(selected_slide_counter, selected_slides_count, "processing", {"slide_index": idx + 1})
        
        has_img, image = has_image_on_left(slide)
        if not has_img:
            stats["skipped_slides"] += 1
            if progress_callback:
                await progress_callback(selected_slide_counter, selected_slides_count, "skipped", {"slide_index": idx + 1})
            continue
        
        # Create new slide
        blank_layout = new_prs.slide_layouts[6]
        new_slide = new_prs.slides.add_slide(blank_layout)
        
        # Copy image
        image_bounds = None
        for shape in slide.shapes:
            if shape.shape_type == 13 and shape.left < prs.slide_width * 0.5:
                image_stream = io.BytesIO(shape.image.blob)
                new_slide.shapes.add_picture(
                    image_stream,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
                image_bounds = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                break
        
        # Translate
        try:
            # Check if cancelled before starting this slide
            if progress_callback:
                await progress_callback(selected_slide_counter, selected_slides_count, "processing", {"phase": "ocr", "slide_index": idx + 1, "slide_total": selected_slides_count})
            
            result = await process_ppt_slide(
                slide, source_lang, target_lang,
                provider, api_key, claude_api_key, model,
                base_font_size, title_adjustment, subject_adjustment,
                log_file=log_path, slide_num=idx + 1
            )
            
            if result and result.get("structure"):
                if image_bounds:
                    margin = Inches(0.2)
                    left_text = image_bounds["left"] + image_bounds["width"] + margin
                    top_text = image_bounds["top"]
                    width_text = new_prs.slide_width - left_text - margin
                    height_text = image_bounds["height"]
                    if width_text < Inches(1.5):
                        left_text = new_prs.slide_width * 0.52
                        top_text = Inches(0.4)
                        width_text = new_prs.slide_width * 0.46
                        height_text = new_prs.slide_height - Inches(0.8)
                else:
                    left_text = new_prs.slide_width * 0.52
                    top_text = Inches(0.4)
                    width_text = new_prs.slide_width * 0.46
                    height_text = new_prs.slide_height - Inches(0.8)
                
                if progress_callback:
                    await progress_callback(selected_slide_counter, selected_slides_count, "processing", {"phase": "render", "slide_index": idx + 1, "slide_total": selected_slides_count})
                apply_formatting_from_structure(
                    new_slide, result["structure"],
                    left_text, top_text, width_text, height_text,
                    base_font_size, title_adjustment, subject_adjustment,
                    ocr_mode=(provider == "ocr_free")
                )
                
                # Manually reduce font size if text overflows the textbox
                # (TEXT_TO_FIT_SHAPE doesn't always reduce enough)
                reduce_font_if_overflow(new_slide, max_iterations=20, min_font_size=6)
                
                stats["processed_slides"] += 1
                
                # Accumulate usage and cost
                if result.get("usage"):
                    stats["total_input_tokens"] += result["usage"].get("input_tokens", 0)
                    stats["total_output_tokens"] += result["usage"].get("output_tokens", 0)
                    
                    # Calculate cost based on provider
                    if provider == "claude":
                        # Claude pricing (per million tokens)
                        if model == "claude-3-haiku-20240307":
                            input_price = 1.0  # $1 per 1M tokens
                            output_price = 5.0  # $5 per 1M tokens
                        else:  # claude-sonnet-4-20250514
                            input_price = 3.0  # $3 per 1M tokens
                            output_price = 15.0  # $15 per 1M tokens
                        
                        slide_cost = (result["usage"]["input_tokens"] / 1_000_000 * input_price) + \
                                    (result["usage"]["output_tokens"] / 1_000_000 * output_price)
                        stats["total_cost"] += slide_cost

                # Track translation method for ALL slides (not just ocr_free)
                slide_info = {
                    "slide": idx + 1,
                    "method": None,
                    "model": None,
                    "error": None
                }
                
                if provider == "ocr_free":
                    slide_info["method"] = result.get("translation_method", "unknown")
                    slide_info["model"] = result.get("translation_model")
                    slide_info["error"] = result.get("openrouter_error")
                elif provider == "claude":
                    slide_info["method"] = "claude_vision"
                    slide_info["model"] = model
                elif provider == "openrouter-vision":
                    slide_info["method"] = "openrouter_vision"
                    slide_info["model"] = model
                
                stats["slide_methods"].append(slide_info)
                
                # Notify backend for progress tracking (ocr_free only needs this for real-time updates)
                if provider == "ocr_free" and progress_callback:
                    await progress_callback(
                        idx + 1, stats["total_slides"], "processing",
                        {"phase": "slide_method", "slide_methods": stats["slide_methods"]}
                    )
                
                if progress_callback:
                    await progress_callback(selected_slide_counter, selected_slides_count, "completed", {"slide_index": idx + 1})
            else:
                # Check if result is an error dict (API error with message)
                if isinstance(result, dict) and "error" in result:
                    error_msg = result["error"]
                else:
                    error_msg = "No result from API (unknown error)"
                print(f"❌ Slide {idx + 1} failed: {error_msg}")
                stats["failed_slides"] += 1
                stats["warnings"].append(f"Slide {idx + 1}: {error_msg}")
                # Add failed slide to methods tracking
                stats["slide_methods"].append({
                    "slide": idx + 1,
                    "method": "unknown",
                    "model": None,
                    "error": error_msg
                })
                if progress_callback:
                    await progress_callback(selected_slide_counter, selected_slides_count, "failed", {"slide_index": idx + 1, "error": error_msg})
        
        except Exception as e:
            # Re-raise if this is a cancellation (don't treat as slide error)
            if "cancelled by user" in str(e).lower():
                raise
            
            print(f"❌ Exception on slide {idx + 1}: {e}")
            import traceback
            traceback.print_exc()
            stats["failed_slides"] += 1
            error_str = str(e)
            stats["warnings"].append(f"Slide {idx + 1}: {error_str}")
            # Add exception to methods tracking
            stats["slide_methods"].append({
                "slide": idx + 1,
                "method": "unknown",
                "model": None,
                "error": error_str
            })
            if progress_callback:
                await progress_callback(selected_slide_counter, selected_slides_count, "failed", {"slide_index": idx + 1, "error": error_str})
    
    # Save
    output = io.BytesIO()
    new_prs.save(output)
    output.seek(0)
    
    # Calculate elapsed time
    stats["elapsed_seconds"] = int(time.time() - stats["start_time"])
    del stats["start_time"]  # Don't send timestamp to frontend, just elapsed time
    
    return output.getvalue(), stats